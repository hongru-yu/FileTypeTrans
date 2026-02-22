"""
PPT文档转换器模块
"""
from pathlib import Path
from typing import Set

from pptx import Presentation


class PptxConverter:
    """PPT文档转换器"""

    # 支持的文件扩展名
    SUPPORTED_EXTENSIONS: Set[str] = {".pptx", ".ppt"}

    def __init__(self):
        """初始化转换器"""
        pass

    def supports(self, file_extension: str) -> bool:
        """
        判断是否支持该文件类型

        Args:
            file_extension: 文件扩展名

        Returns:
            是否支持
        """
        return file_extension.lower() in self.SUPPORTED_EXTENSIONS

    def convert(self, file_path: Path) -> str:
        """
        转换PPT文档为Markdown格式

        Args:
            file_path: 文件路径

        Returns:
            Markdown内容
        """
        # 读取PPT文档
        presentation = Presentation(file_path)

        # 构建Markdown内容
        markdown = []
        markdown.append(f"# 原始文件: {file_path.name}\n")
        markdown.append("---\n")

        # 提取所有幻灯片的文本
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text:
                        markdown.append(text)

        return "\n".join(markdown)
