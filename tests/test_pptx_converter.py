"""
PPT转换器模块的单元测试
"""
import pytest
from pathlib import Path
from unittest.mock import Mock, patch, MagicMock
from src.core.converters.pptx_converter import PptxConverter


class TestPptxConverter:
    """测试PPT转换器类"""

    def test_supports_pptx(self):
        """测试是否支持 .pptx 文件"""
        converter = PptxConverter()
        assert converter.supports(".pptx")

    def test_supports_ppt(self):
        """测试是否支持 .ppt 文件"""
        converter = PptxConverter()
        assert converter.supports(".ppt")

    def test_supports_case_insensitive(self):
        """测试文件扩展名是否大小写不敏感"""
        converter = PptxConverter()
        assert converter.supports(".PPTX")
        assert converter.supports(".PPT")
        assert converter.supports(".Pptx")

    def test_does_not_support_other_formats(self):
        """测试不支持其他格式"""
        converter = PptxConverter()
        assert not converter.supports(".pdf")
        assert not converter.supports(".xlsx")
        assert not converter.supports(".txt")

    @patch('src.core.converters.pptx_converter.Presentation')
    def test_convert_pptx_to_markdown(self, mock_presentation_class):
        """测试转换 .pptx 文件为 Markdown"""
        # 创建 mock presentation 对象
        mock_presentation = Mock()
        mock_presentation_class.return_value = mock_presentation

        # 创建 mock 幻灯片
        mock_slide1 = Mock()
        mock_para1 = Mock()
        mock_para1.text = "第一张幻灯片"
        mock_para2 = Mock()
        mock_para2.text = "这是第一张幻灯片的内容"
        mock_slide1.shapes = [mock_para1, mock_para2]

        mock_slide2 = Mock()
        mock_para3 = Mock()
        mock_para3.text = "第二张幻灯片"
        mock_para4 = Mock()
        mock_para4.text = "这是第二张幻灯片的内容"
        mock_slide2.shapes = [mock_para3, mock_para4]

        mock_presentation.slides = [mock_slide1, mock_slide2]

        # 创建测试文件路径
        test_file = Path("/tmp/test.pptx")

        converter = PptxConverter()
        markdown = converter.convert(test_file)

        # 验证输出
        assert "# 原始文件: test.pptx" in markdown
        assert "第一张幻灯片" in markdown
        assert "这是第一张幻灯片的内容" in markdown
        assert "第二张幻灯片" in markdown
        assert "这是第二张幻灯片的内容" in markdown

        # 验证 Presentation 被正确调用
        mock_presentation_class.assert_called_once_with(test_file)

    @patch('src.core.converters.pptx_converter.Presentation')
    def test_convert_empty_presentation(self, mock_presentation_class):
        """测试转换空PPT文件"""
        mock_presentation = Mock()
        mock_presentation_class.return_value = mock_presentation
        mock_presentation.slides = []

        test_file = Path("/tmp/empty.pptx")

        converter = PptxConverter()
        markdown = converter.convert(test_file)

        assert "# 原始文件: empty.pptx" in markdown

    @patch('src.core.converters.pptx_converter.Presentation')
    def test_convert_presentation_with_empty_slides(self, mock_presentation_class):
        """测试转换包含空幻灯片的PPT文件"""
        mock_presentation = Mock()
        mock_presentation_class.return_value = mock_presentation

        mock_slide1 = Mock()
        mock_para1 = Mock()
        mock_para1.text = "第一张"
        mock_slide1.shapes = [mock_para1]

        mock_slide2 = Mock()
        mock_slide2.shapes = []

        mock_slide3 = Mock()
        mock_para2 = Mock()
        mock_para2.text = "第三张"
        mock_slide3.shapes = [mock_para2]

        mock_presentation.slides = [mock_slide1, mock_slide2, mock_slide3]

        test_file = Path("/tmp/mixed.pptx")

        converter = PptxConverter()
        markdown = converter.convert(test_file)

        assert "第一张" in markdown
        assert "第三张" in markdown

    @patch('src.core.converters.pptx_converter.Presentation')
    def test_convert_presentation_with_special_characters(self, mock_presentation_class):
        """测试转换包含特殊字符的PPT文件"""
        mock_presentation = Mock()
        mock_presentation_class.return_value = mock_presentation

        mock_slide = Mock()
        mock_para = Mock()
        mock_para.text = "中文：你好\n数字：123\n符号：@#$%"
        mock_slide.shapes = [mock_para]

        mock_presentation.slides = [mock_slide]

        test_file = Path("/tmp/special.pptx")

        converter = PptxConverter()
        markdown = converter.convert(test_file)

        assert "中文：你好" in markdown
        assert "数字：123" in markdown
        assert "符号：@#$%" in markdown

    @patch('src.core.converters.pptx_converter.Presentation')
    def test_convert_large_presentation(self, mock_presentation_class):
        """测试转换大PPT文件"""
        mock_presentation = Mock()
        mock_presentation_class.return_value = mock_presentation

        # 创建10张幻灯片
        slides = []
        for i in range(10):
            mock_slide = Mock()
            mock_para = Mock()
            mock_para.text = f"幻灯片 {i+1}"
            mock_slide.shapes = [mock_para]
            slides.append(mock_slide)

        mock_presentation.slides = slides

        test_file = Path("/tmp/large.pptx")

        converter = PptxConverter()
        markdown = converter.convert(test_file)

        assert "幻灯片 1" in markdown
        assert "幻灯片 10" in markdown

    @patch('src.core.converters.pptx_converter.Presentation')
    def test_convert_presentation_with_with_multiple_shapes_per_slide(self, mock_presentation_class):
        """测试转换每张幻灯片包含多个形状的PPT文件"""
        mock_presentation = Mock()
        mock_presentation_class.return_value = mock_presentation

        mock_slide = Mock()
        mock_para1 = Mock()
        mock_para1.text = "标题"
        mock_para2 = Mock()
        mock_para2.text = "副标题"
        mock_para3 = Mock()
        mock_para3.text = "正文内容"
        mock_para4 = Mock()
        mock_para4.text = "更多内容"
        mock_slide.shapes = [mock_para1, mock_para2, mock_para3, mock_para4]

        mock_presentation.slides = [mock_slide]

        test_file = Path("/tmp/multishapes.pptx")

        converter = PptxConverter()
        markdown = converter.convert(test_file)

        assert "标题" in markdown
        assert "副标题" in markdown
        assert "正文内容" in markdown
        assert "更多内容" in markdown

    def test_convert_file_not_found(self):
        """测试文件不存在的错误处理"""
        test_file = Path("/tmp/nonexistent.pptx")

        converter = PptxConverter()

        with pytest.raises(Exception):
            converter.convert(test_file)

    @patch('src.core.converters.pptx_converter.Presentation')
    def test_convert_corrupted_presentation(self, mock_presentation_class):
        """测试损坏PPT的错误处理"""
        from pptx.exc import PackageNotFoundError

        mock_presentation_class.side_effect = PackageNotFoundError("Presentation is corrupted")

        test_file = Path("/tmp/corrupted.pptx")

        converter = PptxConverter()

        with pytest.raises(PackageNotFoundError):
            converter.convert(test_file)

    def test_init(self):
        """测试初始化"""
        converter = PptxConverter()
        assert converter is not None
        assert hasattr(converter, 'supports')
        assert hasattr(converter, 'convert')
